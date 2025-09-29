# Initialization and Setup
import speech_recognition as sr
import requests
import json
import pyttsx3
import os
import difflib
import tempfile
import shutil
import re
import time
from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify, send_from_directory
import sqlite3
from datetime import datetime, timedelta
import pytz
from docx import Document
from werkzeug.utils import secure_filename
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import GPT4All
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
import threading
import logging
from PyPDF2 import PdfReader
import uuid
import pptx
import docx
import PyPDF2
from database import init_db
init_db()

app = Flask(__name__)
app.secret_key = 'rishil-secret-key'  
app.permanent_session_lifetime = timedelta(days=30)

def get_db_connection():
    conn = sqlite3.connect('ai_learning_platform.db')
    conn.row_factory = sqlite3.Row
    return conn

# Configuration and Constants
API_KEY = "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"
recognizer = sr.Recognizer()
tts_engine = pyttsx3.init()

PERSONAL_DATA = {
    "name": "My name is TinoBot",
    "creator": "I was created by Tinos Software and Security Solutions.",
    "birthday": "I don't have an official birthday, but I was created to assist and learn!",
    "creation_location": "I was created by Tinos in Kochi, India, who specializes in AI.",
    "languages": "I can understand and communicate in multiple languages, including English, Hindi, Malayalam, Tamil, and more.",
    "features": "I have several advanced AI capabilities: Natural Language Understanding, Knowledge Base, Summarization, Sentiment Analysis, Multilingual Support, and Context Awareness.",
    "feelings": "Thanks for asking—I'm running smoothly and ready to help! 😊.",
    "day": "Thanks for asking! My day’s going great because I get to chat with you. How about yours?",
    "intro": "Hi, how can I help you today?",
    "ident": "I am TinoBot, invented by Tinos Software and Security Solutions."
}

chat_history = []
user_data = {}
DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
os.makedirs(DATA_DIR, exist_ok=True)

# Utility Functions
def get_current_time():
    ist = pytz.timezone('Asia/Kolkata')
    now = datetime.now(ist)
    return f"The current time is {now.strftime('%I:%M %p')}"

def speak(text):
    tts_engine.say(text)
    tts_engine.runAndWait()

def listen():
    with sr.Microphone() as source:
        recognizer.adjust_for_ambient_noise(source, duration=1)
        print("Listening for wake word...")
        try:
            audio = recognizer.listen(source)
            text = recognizer.recognize_google(audio).lower()
            print(f"Recognized text: {text}")
            return text
        except sr.UnknownValueError:
            print("Could not understand audio.")
            return None
        except sr.RequestError:
            print("Speech recognition service is unavailable.")
            return None

def save_user_data_to_file():
    try:
        user_data_file = os.path.join(DATA_DIR, "user_data.json")
        with open(user_data_file, "w") as f:
            json.dump(user_data, f)
    except Exception as e:
        print(f"Error saving user_data: {e}")

def load_user_data_from_file():
    global user_data
    try:
        user_data_file = os.path.join(DATA_DIR, "user_data.json")
        if os.path.exists(user_data_file):
            with open(user_data_file, "r") as f:
                user_data = json.load(f)
        else:
            user_data = {}
    except Exception as e:
        print(f"Error loading user_data: {e}")
        user_data = {}

def save_chat_history_to_file():
    try:
        chat_history_file = os.path.join(DATA_DIR, "chat_history.json")
        with open(chat_history_file, "w") as f:
            json.dump(chat_history, f)
    except Exception as e:
        print(f"Error saving chat_history: {e}")

def load_chat_history_from_file():
    global chat_history
    try:
        chat_history_file = os.path.join(DATA_DIR, "chat_history.json")
        if os.path.exists(chat_history_file):
            with open(chat_history_file, "r") as f:
                chat_history = json.load(f)
        else:
            chat_history = []
    except Exception as e:
        print(f"Error loading chat_history: {e}")
        chat_history = []

def clear_chat_history():
    global chat_history
    chat_history = []
    chat_history_file = os.path.join(DATA_DIR, "chat_history.json")
    if os.path.exists(chat_history_file):
        os.remove(chat_history_file)
    print("Chat history cleared.")

def check_personal_data(query):
    global user_data
    query = query.lower().strip()

    personal_questions = {
        "name": ["what is your name", "your name", "can you tell me your name", "what should i call you", "may i know your name"],
        "creator": ["who created you", "who made you", "who is your creator", "who developed you", "who is behind your creation", "who constructed you", "who is responsible for your creation", "who programmed you"],
        "birthday": ["your birthday", "when is your birthday", "birth date"],
        "creation_location": ["where were you created", "where were you made"],
        "languages": ["what languages do you know", "which languages can you speak"],
        "features": ["what can you do", "tell me your features"],
        "feelings": ["how are you", "is everything alright", "are you ok", "is everything fine for you"],
        "day": ["how is your day"],
        "intro": ["hi", "hlo", "hy", "hello", "hai"],
        "ident": ["who are you"],
        "time": ["what time is it", "what's the time", "tell me the time", "current time", "what is the current time", "tell me the current time", "what is the time now"],
        "user_name_query": ["what is my name", "what's my name", "do you know my name", "tell me my name"],
        "user_name_set": ["my name is", "i am", "call me"]
    }

    for phrase in personal_questions["user_name_set"]:
        if query.startswith(phrase):
            name = query.replace(phrase, "").strip()
            if name:
                user_data["name"] = name.capitalize()
                save_user_data_to_file()
                return f"Nice to meet you, {user_data['name']}! I'll remember your name."
            return "Sorry, I couldn't understand your name. Please try again by saying 'My name is [your name]'."

    for phrase in personal_questions["user_name_query"]:
        if query == phrase:
            if "name" in user_data:
                return f"Your name is {user_data['name']}."
            return "I don't know your name yet. Please tell me your name by saying 'My name is [your name]'."

    for key, phrases in personal_questions.items():
        if key in ["user_name_query", "user_name_set"]:
            continue
        if query in phrases:
            if key == "time":
                return get_current_time()
            return PERSONAL_DATA[key]


def is_study_related(query):
    """Allow only study-related topics"""
    study_keywords = {
        'math', 'mathematics', 'algebra', 'calculus', 'geometry', 'trigonometry',
        'science', 'physics', 'chemistry', 'biology', 'astronomy', 'geology',
        'history', 'geography', 'civics', 'economics', 'political science',
        'language', 'english', 'literature', 'grammar', 'writing', 'essay',
        'programming', 'computer science', 'coding', 'algorithm', 'data structure',
        'engineering', 'medicine', 'law', 'business', 'accounting', 'finance',
        'philosophy', 'psychology', 'sociology', 'art history', 'architecture',
        'study', 'learn', 'education', 'homework', 'assignment', 'exam', 'test',
        'research', 'thesis', 'dissertation', 'academic', 'school', 'college',
        'university', 'course', 'lecture', 'tutorial', 'textbook', 'curriculum',
        'quiz', 'assessment', 'project', 'presentation', 'seminar', 'workshop',
        'study tips', 'study skills', 'revision', 'notes', 'study guide', 'flashcards',
        'study techniques', 'study methods', 'study strategies', 'study habits',
        'study plan', 'study schedule', 'study resources', 'study materials',
        'exam preparation', 'test preparation', 'academic writing', 'research paper',
        'thesis writing', 'dissertation writing', 'academic research', 'literature review',
        'experiment', 'lab report', 'field study', 'case study', 'data analysis',
        'exploration', 'hypothesis', 'theory', 'concept', 'principle', 'law',
        'definition', 'explanation', 'example', 'problem solving', 'critical thinking',
        'analytical skills', 'logical reasoning', 'creative thinking', 'innovation',
        'problem-based learning', 'project-based learning', 'active learning',
        'collaborative learning', 'peer learning', 'self-directed learning', 'lifelong learning'
        'technology', 'AI', 'artificial intelligence', 'machine learning', 'deep learning',
        'data science', 'big data', 'cloud computing', 'cybersecurity', 'networking',
        'software development', 'web development', 'mobile app development', 'game development',
        'database management', 'information systems', 'IT', 'internet of things',
        'blockchain', 'cryptocurrency', 'virtual reality', 'augmented reality', 'mixed reality',
        'quantum computing', 'natural language processing', 'computer vision', 'robotics',
        'automation', 'digital marketing', 'SEO', 'content marketing', 'social media',
        'e-commerce', 'online learning', 'edtech', 'distance education', 'virtual classroom',
        'online course', 'MOOC', 'learning management system', 'LMS', 'educational technology',
        'academic integrity', 'plagiarism', 'citation', 'referencing', 'bibliography',
        'APA', 'MLA', 'Chicago style', 'Harvard referencing', 'IEEE citation',
        'study group', 'tutoring', 'mentoring', 'academic advising', 'career counseling',
        'college admissions', 'university applications', 'scholarships', 'grants',
        'fellowships', 'internships', 'co-op programs', 'study abroad', 'exchange programs',
        'academic conferences', 'workshops', 'seminars', 'webinars', 'online forums',
        'academic journals', 'research articles', 'peer-reviewed papers', 'conference proceedings',
        'academic books', 'textbooks', 'reference books', 'encyclopedias', 'dictionaries',
        'academic databases', 'library resources', 'digital libraries', 'open access journals',
        'academic networking', 'professional associations', 'academic societies',
        'academic conferences', 'research grants', 'funding opportunities', 'academic publishing',
        'academic careers', 'teaching', 'lecturing', 'research', 'academic administration',
        'academic leadership', 'academic governance', 'academic policy', 'academic standards',
        'academic freedom', 'academic ethics', 'academic misconduct', 'academic discipline',
        'curriculum development', 'course design', 'instructional design',
        'assessment and evaluation', 'learning outcomes', 'competency-based education',
        'blended learning', 'flipped classroom', 'hybrid learning', 'personalized learning',
        'gamification', 'simulation', 'role-playing', 'case-based learning', 'problem-based learning',
        'project-based learning', 'service learning', 'experiential learning', 'fieldwork',
        'internship', 'practicum', 'cooperative education', 'apprenticeship', 'mentorship',
        'academic support', 'tutoring services', 'writing center', 'math lab', 'science lab',
        'library services', 'academic advising', 'career services', 'counseling services', 'define',
        'explain', 'describe', 'summarize', 'compare', 'contrast', 'analyze', 'evaluate',
        'synthesize', 'apply', 'interpret', 'illustrate', 'demonstrate', 'justify', 'what',
    }

    query_lower = query.lower()
    # Check for single-word OR phrase matches
    return (
        any(word in study_keywords for word in query_lower.split()) or  # Single-word
        any(phrase in query_lower for phrase in study_keywords if ' ' in phrase)  # Phrase
    )

# Chat Functionality
def query_deepseek(prompt):
    global chat_history
    cleaned = prompt.strip().lower().replace("  ", " ")
    if cleaned in ["clear history", "clearhistory"]:
        clear_chat_history()
        return "Chat history has been cleared."

    personal_response = check_personal_data(prompt)
    if personal_response:
        chat_history.append({"role": "user", "content": prompt})
        chat_history.append({"role": "assistant", "content": personal_response})
        save_chat_history_to_file()
        return personal_response
    
    # Check if question is study-related
    if not is_study_related(prompt):
        response = "I'm sorry, but I'm designed to only answer study-related questions. Please ask about academic subjects."
        chat_history.append({"role": "user", "content": prompt})
        chat_history.append({"role": "assistant", "content": response})
        save_chat_history_to_file()
        return response

    max_history = 5
    relevant_history = chat_history[-max_history*2:] if len(chat_history) > max_history*2 else chat_history
    messages = [{"role": "system", "content": "I am a helpful AI assistant for academic learning. Provide concise, accurate answers to study questions."}] + relevant_history + [{"role": "user", "content": prompt}]
    
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    json_data = {"model": "deepseek/deepseek-chat-v3.1:free", "messages": messages, "temperature": 0.7}


    try:
        response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=json_data)
        response.raise_for_status()
    except requests.RequestException as e:
        if response.status_code == 429:
            return "API Error: Quota exceeded. Try again later or upgrade your plan."
        return f"API Request Error: {e}"

    if response.status_code == 200:
        try:
            result = response.json()
            if "choices" not in result:
                return f"API Error: Missing 'choices' key. Response: {result}"
            assistant_response = result["choices"][0]["message"]["content"]
            chat_history.append({"role": "user", "content": prompt})
            chat_history.append({"role": "assistant", "content": assistant_response})
            save_chat_history_to_file()
            return assistant_response
        except (KeyError, IndexError, json.JSONDecodeError) as e:
            return f"API Response Error: {e}. Response: {response.text}"
    return f"API Error: {response.status_code}"


# document upload

# Add these configurations at the top of your file
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'uploads')
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'}
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/upload', methods=['POST'])
def upload_document():
    if 'file' not in request.files:
        return jsonify({'success': False, 'message': 'No file part'})
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'success': False, 'message': 'No selected file'})
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)
        
        # Here you would process the file (extract text, etc.)
        # For now, we'll just return success
        return jsonify({
            'success': True,
            'message': 'File uploaded successfully',
            'filename': filename
        })
    
    return jsonify({'success': False, 'message': 'File type not allowed'})

@app.route('/process_document', methods=['POST'])
def process_document():
    data = request.get_json()
    filename = data.get('filename')
    
    if not filename:
        return jsonify({'success': False, 'message': 'Filename not provided'})
    
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    
    if not os.path.exists(filepath):
        return jsonify({'success': False, 'message': 'File not found'})
    
    try:
        # Extract text from the document
        text = extract_text_from_file(filepath)
        
        # Process the text (you might want to summarize, answer questions, etc.)
        response = query_deepseek(f"Please analyze this document and be ready to answer questions about it: {text[:2000]}... [document truncated]")
        
        return jsonify({
            'success': True,
            'message': 'Document processed successfully',
            'summary': response
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'Error processing document: {str(e)}'
        })

def extract_text_from_file(filepath):
    # Implement text extraction based on file type
    if filepath.endswith('.pdf'):
        return extract_text_from_pdf(filepath)
    elif filepath.endswith(('.doc', '.docx')):
        return extract_text_from_doc(filepath)
    elif filepath.endswith(('.ppt', '.pptx')):
        return extract_text_from_ppt(filepath)
    else:  # txt file
        with open(filepath, 'r') as f:
            return f.read()

def extract_text_from_pdf(filepath):
    try:
        import PyPDF2
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
    except ImportError:
        return "PDF processing requires PyPDF2. Please install it."

def extract_text_from_doc(filepath):
    try:
        import docx
        doc = docx.Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except ImportError:
        return "DOCX processing requires python-docx. Please install it."

def extract_text_from_ppt(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except ImportError:
        return "PPTX processing requires python-pptx. Please install it."
    

# Web Routes
@app.route('/')
def index():
    return redirect(url_for('home'))

@app.route('/home')
def home():
    return render_template('home.html')

@app.route('/login')
def login_page():
    return render_template('login.html')

@app.route('/login', methods=['POST'])
def login():
    user_type = request.form.get('user_type')
    user_id = request.form.get('user_id')
    name = request.form.get('name')

    conn = get_db_connection()
    cursor = conn.cursor()

    if user_type == 'student':
        cursor.execute("SELECT * FROM students WHERE student_id = ? AND name = ?", (user_id, name))
    elif user_type == 'teacher':
        cursor.execute("SELECT * FROM teachers WHERE teacher_id = ? AND name = ?", (user_id, name))
    else:
        flash("Invalid user type.", "danger")
        return redirect(url_for('index'))

    user = cursor.fetchone()
    conn.close()

    if not user:
        flash("Invalid ID or Name.", "danger")
        return redirect(url_for('index'))

    session.permanent = True
    session['user_type'] = user_type
    session['user_id'] = user_id
    session['name'] = name
    session['chat_history'] = []

    return redirect(url_for('dashboard'))

@app.route('/dashboard')
def dashboard():
    if 'user_id' not in session:
        return redirect(url_for('index'))
    
    conn = get_db_connection()
    
    try:
        if session['user_type'] == 'teacher':
            # For teachers, get all students and their progress
            students = conn.execute('''
                SELECT s.student_id, s.name, 
                       COUNT(p.id) as quiz_count,
                       AVG(p.quiz_score * 100.0 / p.max_score) as avg_score
                FROM students s
                LEFT JOIN progress p ON s.student_id = p.student_id
                GROUP BY s.student_id, s.name
                ORDER BY s.student_id
            ''').fetchall()
        else:
            # For students, get their own progress and session activity
            students = None
            progress = conn.execute('''
                SELECT topic, quiz_score, max_score, timestamp
                FROM progress
                WHERE student_id = ?
                ORDER BY timestamp DESC
                LIMIT 5
            ''', (session['user_id'],)).fetchall()
            
            activity = conn.execute('''
                SELECT login_time, logout_time, duration_minutes
                FROM session_activity
                WHERE student_id = ?
                ORDER BY login_time DESC
                LIMIT 5
            ''', (session['user_id'],)).fetchall()
            
    except sqlite3.OperationalError as e:
        flash(f"Database error: {str(e)}", 'error')
        students = []
        progress = []
        activity = []
    
    conn.close()
    
    if session['user_type'] == 'teacher':
        return render_template('teacher.html', 
                             teacher=session['name'], 
                             students=students)
    else:
        return render_template('dashboard.html', 
                             name=session['name'], 
                             user_type=session['user_type'],
                             progress=progress,
                             activity=activity)


@app.route('/students')
def list_students():
    if 'user_id' not in session or session['user_type'] != 'teacher':
        return redirect(url_for('index'))
    
    conn = get_db_connection()
    students = conn.execute('''
        SELECT student_id, name, created_at
        FROM students
        ORDER BY created_at DESC
    ''').fetchall()
    
    conn.close()
    return render_template('add_student.html', students=students)

@app.route('/save_student', methods=['POST'])
def save_student():
    if 'user_id' not in session or session['user_type'] != 'teacher':
        return redirect(url_for('index'))
    
    student_id = request.form.get('student_id', '').strip()
    name = request.form.get('name', '').strip()
    
    if not student_id or not name:
        flash('Both student ID and name are required', 'error')
        return redirect(url_for('list_students'))
    
    conn = get_db_connection()
    try:
        # Explicitly set creation timestamp
        created_at = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        conn.execute(
            "INSERT INTO students (student_id, name, created_at) VALUES (?, ?, ?)", 
            (student_id, name, created_at)
        )
        conn.commit()
        flash('Student added successfully', 'success')
    except sqlite3.IntegrityError:
        flash('Student ID already exists', 'error')
    except sqlite3.Error as e:
        flash(f'Database error: {str(e)}', 'error')
    finally:
        conn.close()
    
    return redirect(url_for('list_students'))

@app.route('/delete_student/<student_id>', methods=['DELETE'])
def delete_student(student_id):
    if 'user_id' not in session or session['user_type'] != 'teacher':
        return jsonify({'success': False, 'error': 'Unauthorized'}), 403
    
    conn = get_db_connection()
    try:
        # Verify student exists
        student = conn.execute(
            "SELECT name FROM students WHERE student_id = ?", 
            (student_id,)
        ).fetchone()
        
        if not student:
            return jsonify({'success': False, 'error': 'Student not found'}), 404
        
        # Delete related records
        conn.execute("DELETE FROM progress WHERE student_id = ?", (student_id,))
        conn.execute("DELETE FROM session_activity WHERE user_id = ? AND user_type = 'student'", (student_id,))
        conn.execute("DELETE FROM quiz_files WHERE student_id = ?", (student_id,))
        conn.execute("DELETE FROM user_settings WHERE user_id = ? AND user_type = 'student'", (student_id,))
        
        # Delete the student
        conn.execute("DELETE FROM students WHERE student_id = ?", (student_id,))
        conn.commit()
        
        return jsonify({
            'success': True,
            'message': f"Student {student['name']} deleted successfully"
        })
    except sqlite3.Error as e:
        conn.rollback()
        return jsonify({
            'success': False,
            'error': f"Database error: {str(e)}"
        }), 500
    finally:
        conn.close()

@app.route("/chat", methods=["POST"])
def chat():
    if "user_id" not in session:
        return jsonify({"error": "Not logged in", "status": 403}), 403

    message = request.json.get("message", "")
    if isinstance(message, list):
        message = " ".join(message)
    message = message.strip()

    if not message:
        return jsonify({"error": "Empty message", "status": 400}), 400

    session.setdefault("chat_history", [])

    if message.lower() == "clear history":
        session["chat_history"] = []
        return jsonify({"response": "Chat history has been cleared.", "status": 200})

    personal_response = check_personal_data(message)
    if personal_response:
        session["chat_history"].extend([{"role": "user", "content": message}, {"role": "assistant", "content": personal_response}])
        return jsonify({"response": personal_response, "status": 200})

    session["chat_history"].append({"role": "user", "content": message})
    response = query_deepseek(message)
    session["chat_history"].append({"role": "assistant", "content": response})
    return jsonify({"response": response, "status": 200})

@app.route('/voice', methods=['POST'])
def voice_input():
    audio = request.files['audio']
    return jsonify({'text': "Voice input not implemented yet"})

@app.route('/logout')
def logout():
    session.clear()
    flash("Logged out successfully.", "info")
    return redirect(url_for('index'))


# Admin Routes
from werkzeug.security import generate_password_hash, check_password_hash
from functools import wraps

import sqlite3
def get_db_connection():
    conn = sqlite3.connect('ai_learning_platform.db')
    conn.row_factory = sqlite3.Row
    return conn

@app.route('/admin/login', methods=['GET', 'POST'])
def admin_login():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        
        conn = get_db_connection()
        admin = conn.execute(
            'SELECT * FROM admins WHERE email = ?', 
            (email,)
        ).fetchone()
        conn.close()
        
        if admin and check_password_hash(admin['password_hash'], password):
            session['user_id'] = admin['admin_id']
            session['user_type'] = 'admin'
            session['name'] = admin['name']
            flash('Logged in successfully as admin', 'success')
            return redirect(url_for('admin_dashboard'))
        
        flash('Invalid email or password', 'error')
    
    return render_template('admin_login.html')



@app.route('/admin/logout')
def admin_logout():
    session.clear()
    flash('Logged out successfully', 'success')
    return redirect(url_for('home'))

@app.route('/admin/dashboard')
def admin_dashboard():
    if 'user_id' not in session or session['user_type'] != 'admin':
        return redirect(url_for('admin_login'))
    
    conn = get_db_connection()
    teachers = conn.execute('''
        SELECT t.teacher_id, t.name, t.created_at, 
               a.name as created_by_name
        FROM teachers t
        JOIN admins a ON t.created_by = a.admin_id
        ORDER BY t.created_at DESC
    ''').fetchall()
    conn.close()
    
    return render_template('admin_dashboard.html', teachers=teachers)

@app.route('/admin/add_teacher', methods=['POST'])
def add_teacher():
    if 'user_id' not in session or session['user_type'] != 'admin':
        return jsonify({'success': False, 'error': 'Unauthorized'}), 403
    
    teacher_id = request.form.get('teacher_id', '').strip()
    name = request.form.get('name', '').strip()
    
    if not teacher_id or not name:
        return jsonify({'success': False, 'error': 'Both teacher ID and name are required'}), 400
    
    conn = get_db_connection()
    try:
        # Check if teacher ID already exists
        existing = conn.execute(
            "SELECT 1 FROM teachers WHERE teacher_id = ?", 
            (teacher_id,)
        ).fetchone()
        
        if existing:
            return jsonify({'success': False, 'error': 'Teacher ID already exists'}), 400
        
        # Insert new teacher with just ID and name
        conn.execute(
            "INSERT INTO teachers (teacher_id, name, created_by) VALUES (?, ?, ?)", 
            (teacher_id, name, session['user_id'])
        )
        conn.commit()
        
        return jsonify({
            'success': True,
            'message': 'Teacher added successfully',
            'teacher': {
                'teacher_id': teacher_id,
                'name': name,
                'created_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
        })
    except sqlite3.Error as e:
        return jsonify({'success': False, 'error': f'Database error: {str(e)}'}), 500
    finally:
        conn.close()

@app.route('/admin/delete_teacher/<teacher_id>', methods=['DELETE'])
def delete_teacher(teacher_id):
    if 'user_id' not in session or session['user_type'] != 'admin':
        return jsonify({'success': False, 'error': 'Unauthorized'}), 403
    
    conn = get_db_connection()
    try:
        teacher = conn.execute(
            "SELECT name FROM teachers WHERE teacher_id = ?", 
            (teacher_id,)
        ).fetchone()
        
        if not teacher:
            return jsonify({'success': False, 'error': 'Teacher not found'}), 404
        
        # Delete the teacher
        conn.execute("DELETE FROM teachers WHERE teacher_id = ?", (teacher_id,))
        conn.commit()
        
        return jsonify({
            'success': True,
            'message': f"Teacher {teacher['name']} deleted successfully"
        })
    except sqlite3.Error as e:
        return jsonify({
            'success': False,
            'error': f"Database error: {str(e)}"
        }), 500
    finally:
        conn.close()



# Configure file uploads
UPLOAD_CHAT = os.path.join("static", "chat_uploads")
os.makedirs(UPLOAD_CHAT, exist_ok=True)
app.config["UPLOAD_CHAT"] = UPLOAD_CHAT

@app.route("/teacher_chat")
def teacher_chat():
    if "user_id" not in session or session["user_type"] != "teacher":
        return redirect(url_for("index"))
    conn = get_db_connection()
    students = conn.execute("SELECT student_id, name FROM students").fetchall()
    conn.close()
    return render_template("teacher_select_chat.html", students=students)

@app.route("/student_chat")
def student_chat():
    if "user_id" not in session or session["user_type"] != "student":
        return redirect(url_for("index"))
    conn = get_db_connection()
    teachers = conn.execute("SELECT teacher_id, name FROM teachers").fetchall()
    conn.close()
    if not teachers:
        flash("No teachers available.", "warning")
        return redirect(url_for("dashboard"))
    return render_template("student_select_chat.html", teachers=teachers)

@app.route("/teacher_chat/<student_id>")
def teacher_chat_room(student_id):
    if "user_id" not in session or session["user_type"] != "teacher":
        return redirect(url_for("index"))
    conn = get_db_connection()
    student = conn.execute("SELECT name FROM students WHERE student_id = ?", (student_id,)).fetchone()
    conn.close()
    if not student:
        flash("Student not found.", "danger")
        return redirect(url_for("teacher_chat"))
    chat_id = f"{session['user_id']}_{student_id}"
    return render_template("teacher_chat.html", role="teacher", chat_id=chat_id, other_name=student['name'])

@app.route("/student_chat/<teacher_id>")
def student_chat_room(teacher_id):
    if "user_id" not in session or session["user_type"] != "student":
        return redirect(url_for("index"))
    conn = get_db_connection()
    teacher = conn.execute("SELECT name FROM teachers WHERE teacher_id = ?", (teacher_id,)).fetchone()
    conn.close()
    if not teacher:
        flash("Teacher not found.", "danger")
        return redirect(url_for("student_chat"))
    chat_id = f"{teacher_id}_{session['user_id']}"
    return render_template("teacher_chat.html", role="student", chat_id=chat_id, other_name=teacher['name'])

@app.route("/chat_api/<chat_id>", methods=["GET"])
def get_messages(chat_id):
    since = request.args.get("since", 0)
    conn = get_db_connection()
    rows = conn.execute("""
        SELECT cm.*, 
               CASE 
                   WHEN cm.sender_role = 'teacher' THEN t.name
                   WHEN cm.sender_role = 'student' THEN s.name
               END as sender_name
        FROM chat_messages cm
        LEFT JOIN teachers t ON cm.sender_id = t.teacher_id AND cm.sender_role = 'teacher'
        LEFT JOIN students s ON cm.sender_id = s.student_id AND cm.sender_role = 'student'
        WHERE cm.chat_id=? AND cm.id>? AND cm.is_typing = 0 AND cm.is_deleted = 0
        ORDER BY cm.id
    """, (chat_id, since)).fetchall()

    typing = conn.execute("""
        SELECT cm.sender_role, 
               CASE 
                   WHEN cm.sender_role = 'teacher' THEN t.name
                   WHEN cm.sender_role = 'student' THEN s.name
               END as sender_name
        FROM chat_messages cm
        LEFT JOIN teachers t ON cm.sender_id = t.teacher_id AND cm.sender_role = 'teacher'
        LEFT JOIN students s ON cm.sender_id = s.student_id AND cm.sender_role = 'student'
        WHERE cm.chat_id=? AND cm.is_typing = 1
        ORDER BY cm.id DESC
        LIMIT 1
    """, (chat_id,)).fetchone()
    conn.close()
    return jsonify({
        "messages": [dict(r) for r in rows],
        "typing": dict(typing) if typing else None
    })

@app.route("/chat_api/<chat_id>", methods=["POST"])
def post_message(chat_id):
    if "user_type" not in session or "user_id" not in session:
        return jsonify({"error": "Unauthorized"}), 401
    text = request.form.get("text", "").strip()
    sender_role = session["user_type"]
    sender_id = session["user_id"]
    file_url = filename = file_size = None
    duration = request.form.get("duration", "0:00")

    if not text and "file" not in request.files:
        return jsonify({"error": "Empty message"}), 400

    if "file" in request.files:
        f = request.files["file"]
        if f.filename:
            try:
                filename = secure_filename(f.filename)
                ext = os.path.splitext(filename)[1].lower()
                safe_name = str(uuid.uuid4()) + ext
                path = os.path.join(app.config["UPLOAD_CHAT"], safe_name)
                os.makedirs(os.path.dirname(path), exist_ok=True)
                f.save(path)
                file_url = url_for("download_chat", filename=safe_name, _external=False)
                file_size = os.path.getsize(path)
                if ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp']:
                    file_url = url_for('static', filename=f"chat_uploads/{safe_name}", _external=False)
            except Exception as e:
                return jsonify({"error": f"File upload failed: {str(e)}"}), 500

    try:
        ist = pytz.timezone('Asia/Kolkata')
        current_time = datetime.now(ist)
    except Exception as e:
        current_time = datetime.utcnow()
        print(f"Timezone error, falling back to UTC: {str(e)}")

    conn = None
    try:
        conn = get_db_connection()
        conn.execute("""
            INSERT INTO chat_messages 
            (chat_id, sender_role, sender_id, text, file_url, filename, 
             duration, file_size, created_at, is_read)
            VALUES (?,?,?,?,?,?,?,?,?,?)
        """, (chat_id, sender_role, sender_id, text, file_url, filename, 
              duration, file_size, current_time, 0))
        conn.commit()
        return "", 204
    except sqlite3.Error as e:
        return jsonify({"error": f"Database error: {str(e)}"}), 500
    finally:
        if conn:
            conn.close()

@app.route("/chat_api/<chat_id>/clear", methods=["POST"])
def clear_chat(chat_id):
    if "user_type" not in session or "user_id" not in session:
        return jsonify({"error": "Unauthorized"}), 401
    try:
        teacher_id, student_id = chat_id.split('_')
        if session['user_type'] == 'teacher' and session['user_id'] != teacher_id:
            return jsonify({"error": "Unauthorized"}), 403
        if session['user_type'] == 'student' and session['user_id'] != student_id:
            return jsonify({"error": "Unauthorized"}), 403
        conn = get_db_connection()
        deleted_count = conn.execute("""
            UPDATE chat_messages 
            SET is_deleted = 1,
                deleted_by = ?,
                deleted_at = CURRENT_TIMESTAMP
            WHERE chat_id = ? AND is_deleted = 0
        """, (f"{session['user_type']}:{session['user_id']}", chat_id)).rowcount
        conn.commit()
        conn.execute("""
            DELETE FROM chat_messages 
            WHERE chat_id=? AND is_typing=1
        """, (chat_id,))
        conn.commit()
        return jsonify({"success": True, "deleted": deleted_count})
    except Exception as e:
        return jsonify({"error": f"Error clearing chat: {str(e)}"}), 500
    finally:
        if conn:
            conn.close()

@app.route("/download_chat/<path:filename>")
def download_chat(filename):
    return send_from_directory(
        app.config["UPLOAD_CHAT"],
        filename,
        as_attachment=True,
        download_name=request.args.get("name", filename)
    )

@app.route('/teacher_background')
def teacher_background():
    return render_template("teacher.html")


# research finder


@app.route('/research', methods=['GET', 'POST'])
def research():
    if request.method == 'POST':
        query = request.form['query']
        max_retries = 5
        retry_delay = 1
        
        for attempt in range(max_retries):
            try:
                # First try to get English results
                results = search_semantic_scholar(query, prefer_english=True)
                
                # If no English results found, try again without language preference
                if not results:
                    results = search_semantic_scholar(query, prefer_english=False)
                
                if results:
                    break
                
                time.sleep(retry_delay * (attempt + 1))
            except Exception:
                time.sleep(retry_delay * (attempt + 1))
                results = []
        
        if not results:
            flash("No results found after several attempts. Please try different keywords or try again later.", "info")
        
        return render_template('research.html', results=results, query=query)
    
    return render_template('research.html')

def search_semantic_scholar(query, limit=20, prefer_english=True):
    url = "https://api.semanticscholar.org/graph/v1/paper/search"
    params = {
        "query": query,
        "limit": limit * 3,  # Get more results to increase chance of finding English papers
        "fields": "title,abstract,authors,year,url,openAccessPdf",
    }
    
    # If we prefer English, add journal articles filter (more likely to be English)
    if prefer_english:
        params["publicationTypes"] = "JournalArticle"
    
    response = requests.get(url, params=params, timeout=10)
    response.raise_for_status()
    data = response.json()

    results = []
    for paper in data.get("data", []):
        pdf_link = paper["openAccessPdf"]["url"] if paper.get("openAccessPdf") else ""
        
        paper_data = {
            "title": paper.get("title", "No Title"),
            "abstract": paper.get("abstract", "No Abstract"),
            "year": paper.get("year", "Unknown"),
            "authors": ", ".join([a["name"] for a in paper.get("authors", [])]),
            "link": paper.get("url", ""),
            "pdf_link": pdf_link,
            "is_english": False,  # Will be determined below
        }
        
        # Simple English detection (can be enhanced with langdetect if needed)
        text = f"{paper_data['title']} {paper_data['abstract']}".lower()
        common_english_words = {'the', 'and', 'this', 'that', 'with', 'for', 'are', 'was'}
        if any(word in text for word in common_english_words):
            paper_data["is_english"] = True
        
        results.append(paper_data)
    
    # If we prefer English, filter and return English papers first
    if prefer_english:
        english_papers = [p for p in results if p["is_english"]]
        if english_papers:
            return english_papers[:limit]  # Return only up to the original limit
    
    # Fallback: return all papers (either because prefer_english=False or no English papers found)
    return results[:limit]



# Summarizer
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
model_path = "./mistral-7b-instruct-v0.2.Q4_0.gguf"
if not os.path.exists(model_path):
    raise FileNotFoundError(f"Model file not found at {model_path}")

llm = GPT4All(model=model_path, device='cpu', allow_download=False,n_threads=4) # ,max_tokens=3000
prompt_template = """Write a summary of the following text, using only the provided content and excluding any external or unrelated information. The summary must contain exactly {num_sentences} complete sentences. Each sentence must be concise, clear, and end with proper punctuation:
"{context}"
CONCISE SUMMARY:"""
prompt = PromptTemplate(template=prompt_template, input_variables=["context", "num_sentences"])
llm_chain = LLMChain(llm=llm, prompt=prompt)

def summarize_text(input_text, num_sentences=3):
    try:
        # Prevent empty input
        if not input_text.strip():
            return None, "Error: Input text is empty."
        
        # Create a safe temp file path
        temp_path = os.path.join(tempfile.gettempdir(), "summarizer_input.txt")
        
        # Write input text to the file
        with open(temp_path, "w", encoding="utf-8") as f:
            f.write(input_text)
        
        # Load text as document
        loader = TextLoader(temp_path, encoding="utf-8")
        documents = loader.load()
        
        # Clean up temp file
        os.remove(temp_path)
        
        # Run summarization
        return run_summarization(documents, num_sentences)
    
    except Exception as e:
        return None, f"Error: {str(e)}"


def summarize_file(file_path, num_sentences=3):
    try:
        loader = PyPDFLoader(file_path) if file_path.endswith('.pdf') else TextLoader(file_path)
        documents = loader.load()
        return run_summarization(documents, num_sentences)
    except Exception as e:
        return None, f"Error: {str(e)}"

def run_summarization(documents, num_sentences):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    texts = text_splitter.split_documents(documents)
    persist_directory = tempfile.mkdtemp()
    vectordb = Chroma.from_documents(documents=texts, embedding=embeddings, persist_directory=persist_directory)
    retriever = vectordb.as_retriever(search_kwargs={"k": 4})
    retrieved_docs = retriever.invoke("Summarize the content")
    context = "\n".join([doc.page_content for doc in retrieved_docs])
    del vectordb
    for _ in range(3):
        try:
            shutil.rmtree(persist_directory)
            break
        except PermissionError:
            time.sleep(1)
    for _ in range(3):
        try:
            response = llm_chain.invoke({"context": context, "num_sentences": num_sentences})
            output_raw = response["text"]
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', output_raw) if s.strip() and s[-1] in '.!?']
            if len(sentences) >= num_sentences:
                return "\n".join([f"{i+1}. {s}" for i, s in enumerate(sentences[:num_sentences])]), None
        except Exception as e:
            logger.error(f"LLM generation failed: {e}")
            time.sleep(1)
    return None, "Summary generation failed."

@app.route("/summarize", methods=["GET", "POST"])
def summarize():
    summary = ""
    error = None
    if request.method == "POST":
        num_sentences = int(request.form.get("num_sentences", 3))
        text_input = request.form.get("text_input", "").strip()
        file = request.files.get("file_input")
        if file and file.filename:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            summary, error = summarize_file(file_path, num_sentences)
            os.remove(file_path)
        elif text_input:
            summary, error = summarize_text(text_input, num_sentences)
        else:
            error = "Please enter text or upload a file."
    return render_template("summarize.html", summary=summary, error=error)


# Quiz Generator

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx'}
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
QUIZ_API_KEY = "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"


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(filepath):
    ext = filepath.rsplit('.', 1)[-1].lower()
    text = ""
    if ext == "pdf":
        reader = PdfReader(filepath)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif ext == "docx":
        doc = Document(filepath)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif ext == "txt":
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
    return text.strip()

def generate_quiz_from_text(text, api_key, num_questions=5):
    prompt = f"""
Generate exactly {num_questions} multiple-choice questions based on the following text:
{text}

Format each question exactly like this example:
{{
    "question": "What is the capital of France?",
    "options": ["Paris", "London", "Berlin", "Madrid"],
    "answer_index": 0
}}

Return ONLY a valid JSON array of these questions. No additional text or explanations.
"""
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {"model": "google/gemma-3-27b-it:free", "messages": [{"role": "user", "content": prompt}], "temperature": 0.7}

    try:
        response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload)
        response.raise_for_status()
        result = response.json()
        
        content = result["choices"][0]["message"]["content"]
        json_start = content.find('[')
        json_end = content.rfind(']') + 1
        if json_start == -1 or json_end == -1:
            raise ValueError("No valid JSON array found in response")
            
        quiz_data = json.loads(content[json_start:json_end])
        
        # Validate the structure
        if not isinstance(quiz_data, list):
            raise ValueError("Response is not a list")
            
        for question in quiz_data:
            if not all(key in question for key in ['question', 'options', 'answer_index']):
                raise ValueError("Missing required question fields")
            if not isinstance(question['options'], list) or len(question['options']) != 4:
                raise ValueError("Each question must have exactly 4 options")
            if not isinstance(question['answer_index'], int) or question['answer_index'] not in range(4):
                raise ValueError("Invalid answer_index")
        
        return quiz_data
        
    except (requests.RequestException, json.JSONDecodeError, KeyError, ValueError) as e:
        raise Exception(f"Quiz generation failed: {str(e)}")

def generate_quiz_from_topic(topic, api_key, num_questions=5, max_retries=3):
    """Generate quiz questions based on a topic name with automatic retries on failure"""
    prompt = f"""
Generate exactly {num_questions} multiple-choice questions about {topic}. 
Each question must have 4 options with only one correct answer.
Format each question exactly like this example:
{{
    "question": "What is the capital of France?",
    "options": ["Paris", "London", "Berlin", "Madrid"],
    "answer_index": 0
}}

Return ONLY a valid JSON array of these questions. No additional text or explanations.
"""
    
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    
    payload = {
        "model": "google/gemma-3-27b-it:free",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.7,
        "max_tokens": 2000
    }
    
    for attempt in range(max_retries):
        try:
            response = requests.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=30
            )
            response.raise_for_status()
            result = response.json()
            
            content = result["choices"][0]["message"]["content"]
            
            # Clean the response to extract just the JSON
            json_start = content.find('[')
            json_end = content.rfind(']') + 1
            if json_start == -1 or json_end == -1:
                raise ValueError("No valid JSON array found in response")
                
            quiz_data = json.loads(content[json_start:json_end])
            
            # Validate the structure
            if not isinstance(quiz_data, list):
                raise ValueError("Response is not a list")
                
            for question in quiz_data:
                if not all(key in question for key in ['question', 'options', 'answer_index']):
                    raise ValueError("Missing required question fields")
                if not isinstance(question['options'], list) or len(question['options']) != 4:
                    raise ValueError("Each question must have exactly 4 options")
                if not isinstance(question['answer_index'], int) or question['answer_index'] not in range(4):
                    raise ValueError("Invalid answer_index")
            
            return quiz_data
            
        except (requests.RequestException, json.JSONDecodeError, KeyError, ValueError) as e:
            print(f"Attempt {attempt + 1} failed: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(2 ** attempt)  # Exponential backoff
                continue
            else:
                print(f"Max retries ({max_retries}) reached. Giving up.")
                return None

# Quiz Routes
@app.route("/quiz_home", methods=["GET", "POST"])
def quiz_home():
    if request.method == "POST":
        if 'user_type' not in session or session['user_type'] != 'student':
            flash("Please log in as a student to take a quiz.", "warning")
            return redirect(url_for('index'))
        
        topic = request.form["topic"]
        num = int(request.form.get("num", 5))
        try:
            quiz = generate_quiz_from_topic(topic, api_key=QUIZ_API_KEY, num_questions=num)
            if quiz is None:
                flash("Failed to generate quiz. Please try again with a different topic.", "danger")
                return redirect(url_for("quiz_home"))
            
            session["quiz"] = quiz
            session["quiz_topic"] = topic
            return redirect(url_for("take_quiz"))
        except Exception as e:
            flash(f"Error: {str(e)}", "danger")
    return render_template("quiz_home.html")

@app.route("/quiz_upload", methods=["GET", "POST"])
def upload_file():
    if request.method == "POST":
        if 'user_type' not in session or session['user_type'] != 'student':
            flash("Please log in as a student to upload files.", "warning")
            return redirect(url_for('index'))
        
        file = request.files.get("file")
        num = int(request.form.get("num", 5))

        if not file or file.filename == '':
            flash("No file selected.", "danger")
            return redirect(request.url)

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            text = extract_text_from_file(filepath)
            if not text:
                flash("Failed to extract text from file.", "danger")
                return redirect(request.url)

            try:
                quiz = generate_quiz_from_text(text, api_key=QUIZ_API_KEY, num_questions=num)
                session["quiz"] = quiz
                session["quiz_topic"] = f"Document: {filename}"
                return redirect(url_for("take_quiz"))
            except Exception as e:
                flash(f"Quiz generation failed: {str(e)}", "danger")

    return render_template("quiz_upload.html")

@app.route("/quiz", methods=["GET", "POST"])
def take_quiz():
    if 'user_id' not in session or session['user_type'] != 'student' or 'quiz' not in session:
        flash("Invalid quiz session.", "danger")
        return redirect(url_for('index'))
    
    quiz = session.get("quiz")
    if not quiz:
        return redirect(url_for("quiz_home"))

    if request.method == "POST":
        try:
            # Get user answers
            answers = []
            for i in range(len(quiz)):
                answer = request.form.get(f"q{i}")
                answers.append(int(answer) if answer is not None else None)

            # Calculate score
            score = 0
            results = []
            for i, (question, user_answer) in enumerate(zip(quiz, answers)):
                is_correct = user_answer == question['answer_index']
                if is_correct:
                    score += 1
                
                results.append({
                    'question': question['question'],
                    'options': question['options'],
                    'user_answer': user_answer,
                    'user_display': question['options'][user_answer] if user_answer is not None else None,
                    'correct_display': question['options'][question['answer_index']],
                    'is_correct': is_correct
                })

            # Store in database
            conn = get_db_connection()
            conn.execute('''
                INSERT INTO progress (
                    student_id, 
                    quiz_score, 
                    max_score, 
                    topic, 
                    num_questions, 
                    answers,
                    question_details,
                    time_taken
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            ''', (
                session['user_id'],
                score,
                len(quiz),
                session.get("quiz_topic", "General"),
                len(quiz),
                json.dumps(answers),
                json.dumps(results),
                len(quiz) * 30  # 30 seconds per question
            ))
            conn.commit()
            conn.close()

            # Clear the quiz from session
            session.pop('quiz', None)
            session.pop('quiz_topic', None)

            return render_template(
                "quiz_result.html",
                results=results,
                score=score,
                max_score=len(quiz),
                topic=session.get('quiz_topic', 'General'),
                now=datetime.now()
            )

        except Exception as e:
            flash(f"An error occurred: {str(e)}", "danger")
            return redirect(url_for("quiz_home"))

    return render_template("quiz.html", quiz=quiz)




# Progress Tracking
UPLOAD_FOLDER = 'static/uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Progress Tracking Routes
@app.route("/progress")
def progress():
    if "user_type" not in session:
        flash("Please log in first.", "warning")
        return redirect(url_for("index"))
    
    student_id = session["user_id"]
    user_type = session["user_type"]
    conn = get_db_connection()
    cursor = conn.cursor()
    
    if user_type == "teacher":
        cursor.execute("SELECT student_id, name FROM students")
        students = cursor.fetchall()
        conn.close()
        return render_template("teacher.html", students=students)
    
    # Get progress data with proper timezone handling
    cursor.execute("""
        SELECT 
            id,
            quiz_score, 
            max_score, 
            datetime(timestamp, 'localtime') as timestamp,
            topic, 
            num_questions, 
            answers, 
            question_details,
            time_taken
        FROM progress 
        WHERE student_id = ? 
        ORDER BY timestamp DESC
    """, (student_id,))
    
    progress_data = cursor.fetchall()
    
    # Convert to dictionary and format timestamps
    progress_dict = []
    for row in progress_data:
        row_dict = dict(row)
        if row_dict["timestamp"]:
            try:
                # Convert to datetime object and format for display
                dt = datetime.strptime(row_dict["timestamp"], "%Y-%m-%d %H:%M:%S")
                row_dict["timestamp"] = dt.strftime("%Y-%m-%d %I:%M %p")  # Format for display
            except (ValueError, TypeError) as e:
                print(f"Error parsing timestamp: {e}")
                row_dict["timestamp"] = None
        progress_dict.append(row_dict)
    
    # Calculate metrics
    scores = []
    for row in progress_dict:
        if row["max_score"] and row["max_score"] > 0 and row["quiz_score"] is not None:
            scores.append(row["quiz_score"] / row["max_score"] * 100)
    
    average_score = sum(scores) / len(scores) if scores else 0
    
    # Get time spent - updated to use duration_seconds
    cursor.execute("""
        SELECT duration_seconds 
        FROM session_activity 
        WHERE user_id = ? AND logout_time IS NOT NULL AND duration_seconds IS NOT NULL
    """, (student_id,))
    time_data = cursor.fetchall()
    total_time = sum(row["duration_seconds"] or 0 for row in time_data) / 60  # Convert to minutes
    
    # Calculate total question attempts
    question_attempts = sum(row["num_questions"] or 0 for row in progress_dict)
    
    # Generate feedback
    improvement_feedback = generate_feedback(student_id)
    
    conn.close()
    
    # Group progress by topic and extract short topic for display
    topics_progress = {}
    for record in progress_dict:
        topic = record.get("topic", "General") or "General"
        
        # Shorten topic to first sentence or line
        short_topic = topic.split('.')[0] if '.' in topic else topic.split('\n')[0]
        record["short_topic"] = short_topic[:30]  # Limit to 30 characters for safety

        if topic not in topics_progress:
            topics_progress[topic] = []
        topics_progress[topic].append(record)
    
    return render_template("progress.html",
                       progress=progress_dict,
                       topics_progress=topics_progress,
                       feedbacks=improvement_feedback,
                       avg_score=round(average_score, 2),
                       total_time=round(total_time, 2),
                       question_attempts=question_attempts,
                       name=session.get("name", "Student"),
                       user_type=user_type)

@app.route('/view_student_progress', methods=['POST'])
def view_student_progress():
    if "user_type" not in session or session["user_type"] != "teacher":
        flash("Access denied. Only teachers can view student progress.", "danger")
        return redirect(url_for('progress'))
    
    student_id = request.form.get('student_id')
    student_name = request.form.get('student_name')
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Verify student exists
    cursor.execute("""
        SELECT * FROM students 
        WHERE student_id = ? AND name = ?
    """, (student_id, student_name))
    student = cursor.fetchone()
    
    if not student:
        flash("Student not found.", "danger")
        conn.close()
        return redirect(url_for('progress'))
    
    # Get progress data with proper timezone handling
    cursor.execute("""
        SELECT 
            id,
            quiz_score, 
            max_score, 
            datetime(timestamp, 'localtime') as timestamp,
            topic, 
            num_questions, 
            answers, 
            question_details,
            time_taken
        FROM progress 
        WHERE student_id = ? 
        ORDER BY timestamp DESC
    """, (student_id,))
    progress_data = cursor.fetchall()
    
    # Convert to dictionary and format timestamps
    progress_dict = []
    for row in progress_data:
        row_dict = dict(row)
        if row_dict["timestamp"]:
            try:
                dt = datetime.strptime(row_dict["timestamp"], "%Y-%m-%d %H:%M:%S")
                row_dict["timestamp"] = dt.strftime("%Y-%m-%d %I:%M %p")  # Format for display
            except ValueError:
                row_dict["timestamp"] = None
        progress_dict.append(row_dict)
    
    # Calculate metrics
    scores = []
    for row in progress_dict:
        if row["max_score"] and row["max_score"] > 0 and row["quiz_score"] is not None:
            scores.append(row["quiz_score"] / row["max_score"] * 100)
    
    average_score = sum(scores) / len(scores) if scores else 0
    
    # Calculate total time spent - updated to use duration_seconds
    cursor.execute("""
        SELECT duration_seconds 
        FROM session_activity 
        WHERE user_id = ? AND logout_time IS NOT NULL AND duration_seconds IS NOT NULL
    """, (student_id,))
    time_data = cursor.fetchall()
    total_time = sum(row["duration_seconds"] or 0 for row in time_data) / 60  # Convert to minutes
    
    # Calculate total question attempts
    question_attempts = sum(row["num_questions"] or 0 for row in progress_dict)
    
    # Generate feedback
    improvement_feedback = generate_feedback(student_id)
    
    conn.close()
    
    # Group progress by topic and extract short topic for display
    topics_progress = {}
    for record in progress_dict:
        topic = record.get("topic", "General") or "General"
        
        # Shorten topic to first sentence or line
        short_topic = topic.split('.')[0] if '.' in topic else topic.split('\n')[0]
        record["short_topic"] = short_topic[:30]  # Limit to 30 characters for safety

        if topic not in topics_progress:
            topics_progress[topic] = []
        topics_progress[topic].append(record)
    
    return render_template("progress.html",
                       progress=progress_dict,
                       topics_progress=topics_progress,
                       feedbacks=improvement_feedback,
                       avg_score=round(average_score, 2),
                       total_time=round(total_time, 2),
                       question_attempts=question_attempts,
                       name=student_name,
                       user_type=session["user_type"])

@app.route('/clear_progress', methods=['POST'])
def clear_progress():
    if "user_type" not in session or session["user_type"] != "student":
        flash("Unauthorized action.", "danger")
        return redirect(url_for("progress"))

    student_id = session["user_id"]
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Delete progress and session activity
    cursor.execute("DELETE FROM progress WHERE student_id = ?", (student_id,))
    cursor.execute("DELETE FROM session_activity WHERE student_id = ?", (student_id,))
    
    conn.commit()
    conn.close()
    
    flash("Your progress has been cleared successfully.", "success")
    return redirect(url_for("progress"))



def generate_feedback(student_id):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    cursor.execute("""
        SELECT 
            quiz_score, 
            max_score, 
            timestamp,
            topic
        FROM progress 
        WHERE student_id = ? 
        ORDER BY timestamp DESC
    """, (student_id,))
    progress_data = cursor.fetchall()
    conn.close()

    if not progress_data:
        return [
            ["Improvement", "N/A", "No data yet, start with a quiz!"],
            ["Correct Answers", "N/A", "Take a quiz to see results"],
            ["Average Score", "N/A", "Get started to see progress"],
            ["Recommendation", "N/A", "Try a quiz to begin your journey"],
            {"strengths": [], "improvements": []}
        ]

    # Calculate metrics
    total_score = sum(row["quiz_score"] for row in progress_data)
    total_max_score = sum(row["max_score"] for row in progress_data)
    avg_score = (total_score / total_max_score * 100) if total_max_score else 0
    recent_score = progress_data[-1]["quiz_score"] if progress_data else 0
    recent_max = progress_data[-1]["max_score"] if progress_data else 1

    # Improvement calculation
    if len(progress_data) > 1:
        prev_score = progress_data[-2]["quiz_score"]
        improvement = recent_score - prev_score
        improvement_feedback = ["Improvement", f"{improvement}", "Positive" if improvement > 0 else "Negative" if improvement < 0 else "Neutral"]
    else:
        improvement_feedback = ["Improvement", "N/A", "More attempts needed"]

    # Topic analysis
    topic_performance = {}
    for row in progress_data:
        topic = row["topic"] or "General"
        if topic not in topic_performance:
            topic_performance[topic] = {"total": 0, "count": 0}
        if row["max_score"] > 0:
            topic_performance[topic]["total"] += (row["quiz_score"] / row["max_score"]) * 100
            topic_performance[topic]["count"] += 1

    # Generate strengths and weaknesses
    strengths = []
    improvements = []
    
    if topic_performance:
        best_topic = max(topic_performance.items(), key=lambda x: x[1]["total"]/x[1]["count"])
        strengths.append(f"Strong in {best_topic[0]} ({best_topic[1]['total']/best_topic[1]['count']:.1f}%)")
        
        worst_topic = min(topic_performance.items(), key=lambda x: x[1]["total"]/x[1]["count"])
        improvements.append(f"Needs work in {worst_topic[0]} ({worst_topic[1]['total']/worst_topic[1]['count']:.1f}%)")

    if len(progress_data) >= 3:
        last_three = [p["quiz_score"] for p in progress_data[-3:]]
        trend = ((last_three[-1] - last_three[0]) / last_three[0]) * 100 if last_three[0] != 0 else 0
        if trend > 10:
            strengths.append(f"Improving trend (+{trend:.1f}%)")
        elif trend < -10:
            improvements.append(f"Declining trend ({trend:.1f}%)")

    if not strengths:
        if avg_score >= 70:
            strengths.append("Consistently high scores")
        elif avg_score >= 50:
            strengths.append("Average performance")

    if not improvements:
        if avg_score < 50:
            improvements.append("Focus on fundamentals")
        improvements.append("Review incorrect answers")

    return [
        improvement_feedback,
        ["Correct Answers", f"{recent_score}/{recent_max}", "See results"],
        ["Average Score", f"{avg_score:.1f}%", "Good" if avg_score >= 70 else "Needs Work" if avg_score < 50 else "Fair"],
        ["Recommendation", "Practice more in weaker areas", "Actionable step"],
        {"strengths": strengths, "improvements": improvements}
    ]



# Quiz Submission Handler (to be called when quiz is completed)
def save_quiz_result(student_id, score, max_score, topic, answers, question_details, filename=None):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Get current time in IST
    ist = pytz.timezone('Asia/Kolkata')
    current_time = datetime.now(ist)
    
    cursor.execute('''
        INSERT INTO progress (
            student_id, 
            quiz_score, 
            max_score, 
            timestamp,
            topic, 
            num_questions, 
            answers,
            question_details
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
    ''', (
        student_id,
        score,
        max_score,
        current_time.strftime('%Y-%m-%d %H:%M:%S'),
        topic,
        max_score,  # Assuming 1 question per point
        json.dumps(answers),
        json.dumps(question_details)
    ))
    
    if filename:
        cursor.execute('''
            INSERT INTO quiz_files (student_id, filename)
            VALUES (?, ?)
        ''', (student_id, filename))
    
    conn.commit()
    conn.close()

if __name__ == "__main__":
    app.run(debug=True)